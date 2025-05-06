import boto3
import os
from pptx import Presentation
from bedrock import Bedrock  # Assuming Bedrock SDK is installed via requirements.txt

s3 = boto3.client('s3')
bedrock = Bedrock()  # Initialize Bedrock client

def lambda_handler(event, context):
    bucket = event['BUCKET']
    key = event['KEY']
    language = event.get('LANGUAGE', 'en')  # Default to English if not specified
    
    # Download file from S3
    local_file = '/tmp/' + os.path.basename(key)
    s3.download_file(bucket, key, local_file)
    
    # Load PPT
    ppt = Presentation(local_file)
    
    # Extract text and translate (simplified example)
    translated_slides = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                translated_text = bedrock.translate(shape.text, target_language=language)
                shape.text = translated_text  # Update with translated text
                translated_slides.append(slide)
    
    # Save translated PPT
    translated_key = 'translated_' + os.path.basename(key)
    translated_local_file = '/tmp/translated_' + os.path.basename(key)
    ppt.save(translated_local_file)
    
    # Upload translated file back to S3
    s3.upload_file(translated_local_file, bucket, translated_key)
    
    return {
        'statusCode': 200,
        'body': f'Translation complete: {translated_key}'
    }

if __name__ == '__main__':
    # For local testing
    event = {'BUCKET': os.getenv('BUCKET'), 'KEY': os.getenv('KEY'), 'LANGUAGE': os.getenv('LANGUAGE')}
    lambda_handler(event, None)
