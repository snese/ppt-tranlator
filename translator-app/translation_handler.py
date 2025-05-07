import json
import boto3
import os
import argparse
import logging
import time
from botocore.exceptions import ClientError
from tqdm import tqdm
from pptx import Presentation
from pptx.dml.color import RGBColor
from copy import deepcopy

# Set up logging with detailed format
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class BedrockTranslator:
    """
    A translator class using AWS Bedrock with Claude 3.5 Sonnet model for text translation.
    Supports interactive, file-based, and batch translation modes with advanced features.
    """
    
    def __init__(self, region_name="us-west-2"):
        """
        Initialize the translator with AWS Bedrock Runtime client.
        
        Args:
            region_name (str): AWS region name, default is 'us-west-2'.
        """
        self.bedrock_runtime = boto3.client("bedrock-runtime", region_name=region_name)
        self.model_id = "anthropic.claude-3-5-sonnet-20241022-v2:0"  # Using a model ID that works reliably
        self.s3_client = boto3.client('s3', region_name=region_name)
        logger.info(f"Initialized BedrockTranslator with region {region_name}")
    
    def translate(self, text, source_language="auto (en-US)", target_language="zh-TW", 
                  use_reasoning=False, temperature=0.7, max_tokens=3000, top_p=0.9):
        """
        Translate a single piece of text using AWS Bedrock.
        
        Args:
            text (str): Text to translate.
            source_language (str): Source language, default is 'auto (en-US)'.
            target_language (str): Target language, default is 'zh-TW'.
            use_reasoning (bool): Whether to enable extended reasoning for improved accuracy.
            temperature (float): Temperature for model creativity, default is 0.7.
            max_tokens (int): Maximum tokens for response, default is 3000.
            top_p (float): Top P for nucleus sampling, default is 0.9.
            
        Returns:
            str: Translated text.
        """
        max_retries = 3
        for attempt in range(max_retries):
            try:
                # Build translation prompt
                prompt = f"Translate the following text from {source_language} to {target_language}, maintaining the original format, tone, and meaning. Return only the translated text without any additional explanation:\n\n{text}"
                
                if use_reasoning:
                    prompt = f"Translate the following text from {source_language} to {target_language}. First, analyze key terms and style, then provide an accurate translation that preserves the original format, tone, and technical accuracy. Return only the translated text without explanation:\n\n{text}"
                
                # Build messages for Converse API
                messages = [
                    {
                        "role": "user",
                        "content": [{"text": prompt}]
                    }
                ]
                
                # Configure inference parameters
                inference_config = {
                    "maxTokens": max_tokens,
                    "temperature": temperature,
                    "topP": top_p
                }
                
                if use_reasoning:
                    # Enable extended reasoning with token budget
                    inference_config["reasoning"] = {"thinking": {"type": "enabled", "budget": 2000}}
                
                # Call Converse API
                logger.info(f"Invoking Bedrock model {self.model_id} for translation")
                response = self.bedrock_runtime.converse(
                    modelId=self.model_id,
                    messages=messages,
                    inferenceConfig=inference_config
                )
                
                # Extract translated text
                translated_text = response["output"]["message"]["content"][0]["text"]
                logger.info("Translation successful")
                return translated_text
                
            except ClientError as e:
                logger.error(f"Attempt {attempt+1}/{max_retries} failed: {str(e)}")
                if attempt < max_retries - 1:
                    time.sleep(2 ** attempt)  # Exponential backoff
                    continue
                else:
                    raise Exception(f"Failed to translate text after {max_retries} attempts: {str(e)}")
            except Exception as e:
                logger.error(f"Unexpected error in attempt {attempt+1}/{max_retries}: {str(e)}")
                if attempt < max_retries - 1:
                    time.sleep(2 ** attempt)  # Exponential backoff
                    continue
                else:
                    raise Exception(f"Unexpected error during translation after {max_retries} attempts: {str(e)}")
    
    def copy_run_formatting(self, source_run, target_run):
        """
        Copy formatting from source run to target run.
        
        Args:
            source_run: Source run to copy formatting from.
            target_run: Target run to apply formatting to.
        """
        # Copy font properties
        if source_run.font.bold is not None:
            target_run.font.bold = source_run.font.bold
        
        if source_run.font.italic is not None:
            target_run.font.italic = source_run.font.italic
        
        if source_run.font.underline is not None:
            target_run.font.underline = source_run.font.underline
        
        if source_run.font.size is not None:
            target_run.font.size = source_run.font.size
        
        # Copy font color if it exists
        if hasattr(source_run.font.color, 'rgb') and source_run.font.color.rgb is not None:
            target_run.font.color.rgb = source_run.font.color.rgb
        
        # Copy other properties as needed
        if hasattr(source_run.font, 'name') and source_run.font.name is not None:
            target_run.font.name = source_run.font.name
    
    def translate_text_frame(self, text_frame, source_language, target_language, use_reasoning, temperature, max_tokens, top_p):
        """
        Translate text in a text frame while preserving formatting.
        
        Args:
            text_frame: The text frame to translate.
            source_language: Source language.
            target_language: Target language.
            use_reasoning: Whether to enable extended reasoning.
            temperature: Temperature for model creativity.
            max_tokens: Maximum tokens for response.
            top_p: Top P for nucleus sampling.
            
        Returns:
            bool: True if translation is successful, False otherwise.
        """
        if not text_frame.text.strip():
            return True
        
        try:
            # Store original text and formatting information
            original_paragraphs = []
            for p in text_frame.paragraphs:
                p_info = {
                    'text': p.text,
                    'runs': []
                }
                
                # Store alignment and level
                p_info['alignment'] = p.alignment
                p_info['level'] = p.level
                
                # Store run information
                for run in p.runs:
                    run_info = {
                        'text': run.text,
                        'font': {
                            'bold': run.font.bold,
                            'italic': run.font.italic,
                            'underline': run.font.underline,
                            'size': run.font.size,
                            'name': run.font.name if hasattr(run.font, 'name') else None
                        }
                    }
                    
                    # Store color information if available
                    if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                        run_info['font']['color'] = run.font.color.rgb
                    
                    p_info['runs'].append(run_info)
                
                original_paragraphs.append(p_info)
            
            # Translate the entire text
            original_text = text_frame.text
            translated_text = self.translate(
                original_text, source_language, target_language,
                use_reasoning, temperature, max_tokens, top_p
            )
            
            # Split translated text into paragraphs
            translated_paragraphs = translated_text.split('\n')
            
            # Clear the text frame
            for i in range(len(text_frame.paragraphs)):
                p = text_frame.paragraphs[i]
                if i == 0:
                    p.clear()
                else:
                    # Can't remove paragraphs, so clear all but the first
                    for run in p.runs:
                        run.text = ''
            
            # Add translated paragraphs with original formatting
            for i, translated_para_text in enumerate(translated_paragraphs):
                if i == 0:
                    # Use the first paragraph that's already in the text frame
                    p = text_frame.paragraphs[0]
                else:
                    # Add new paragraphs as needed
                    p = text_frame.add_paragraph()
                
                # Apply original paragraph formatting if available
                if i < len(original_paragraphs):
                    if original_paragraphs[i]['alignment'] is not None:
                        p.alignment = original_paragraphs[i]['alignment']
                    if original_paragraphs[i]['level'] is not None:
                        p.level = original_paragraphs[i]['level']
                
                # Add a single run with the translated text
                run = p.add_run()
                run.text = translated_para_text
                
                # Apply formatting from the first run of the original paragraph
                if i < len(original_paragraphs) and original_paragraphs[i]['runs']:
                    original_run_info = original_paragraphs[i]['runs'][0]
                    
                    # Apply font properties
                    if original_run_info['font']['bold'] is not None:
                        run.font.bold = original_run_info['font']['bold']
                    
                    if original_run_info['font']['italic'] is not None:
                        run.font.italic = original_run_info['font']['italic']
                    
                    if original_run_info['font']['underline'] is not None:
                        run.font.underline = original_run_info['font']['underline']
                    
                    if original_run_info['font']['size'] is not None:
                        run.font.size = original_run_info['font']['size']
                    
                    if original_run_info['font']['name'] is not None:
                        run.font.name = original_run_info['font']['name']
                    
                    # Apply color if available
                    if 'color' in original_run_info['font']:
                        run.font.color.rgb = original_run_info['font']['color']
            
            return True
        
        except Exception as e:
            logger.error(f"Error translating text frame: {e}")
            return False
    
    def translate_file(self, input_file, output_file, source_language="auto (en-US)", 
                       target_language="zh-TW", use_reasoning=False, temperature=0.7, 
                       max_tokens=3000, top_p=0.9):
        """
        Translate content from a PowerPoint file and save to a new file.
        
        Args:
            input_file (str): Path to input PowerPoint file.
            output_file (str): Path to save translated PowerPoint file.
            source_language (str): Source language, default is 'auto (en-US)'.
            target_language (str): Target language, default is 'zh-TW'.
            use_reasoning (bool): Whether to enable extended reasoning for improved accuracy.
            temperature (float): Temperature for model creativity, default is 0.7.
            max_tokens (int): Maximum tokens for response, default is 3000.
            top_p (float): Top P for nucleus sampling, default is 0.9.
            
        Returns:
            bool: True if translation is successful, False otherwise.
        """
        try:
            logger.info(f"Processing PowerPoint file: {input_file}")
            
            # Load the presentation
            prs = Presentation(input_file)
            
            # Count total items for progress tracking
            total_shapes = sum(len(slide.shapes) for slide in prs.slides)
            total_notes = sum(1 for slide in prs.slides if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip())
            total_items = total_shapes + total_notes
            processed_items = 0
            
            with tqdm(total=total_items, desc="Translating content") as pbar:
                for slide_idx, slide in enumerate(prs.slides):
                    logger.info(f"Processing slide {slide_idx+1}/{len(prs.slides)}")
                    
                    # Process each shape in the slide
                    for shape_idx, shape in enumerate(slide.shapes):
                        # Handle text in shapes
                        if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                            original_text = shape.text_frame.text
                            logger.info(f"Translating shape text: {original_text[:50]}...")
                            
                            # Translate the text frame with formatting preservation
                            self.translate_text_frame(
                                shape.text_frame, source_language, target_language,
                                use_reasoning, temperature, max_tokens, top_p
                            )
                            
                            logger.info(f"Translated shape text with preserved formatting")
                        
                        # Handle tables explicitly
                        if hasattr(shape, "table"):
                            logger.info(f"Processing table with {len(shape.table.rows)} rows and {len(shape.table.columns)} columns")
                            for row_idx, row in enumerate(shape.table.rows):
                                for col_idx, cell in enumerate(row.cells):
                                    if cell.text_frame.text.strip():
                                        original_cell_text = cell.text_frame.text
                                        logger.info(f"Translating table cell [{row_idx},{col_idx}]: {original_cell_text[:50]}...")
                                        
                                        # Translate the cell text frame with formatting preservation
                                        self.translate_text_frame(
                                            cell.text_frame, source_language, target_language,
                                            use_reasoning, temperature, max_tokens, top_p
                                        )
                                        
                                        logger.info(f"Translated table cell with preserved formatting")
                        
                        processed_items += 1
                        pbar.update(1)
                    
                    # Translate speaker notes if they exist
                    if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
                        notes_slide = slide.notes_slide
                        original_notes = notes_slide.notes_text_frame.text
                        if original_notes.strip():
                            logger.info(f"Translating speaker notes: {original_notes[:50]}...")
                            
                            # Translate the notes text frame with formatting preservation
                            self.translate_text_frame(
                                notes_slide.notes_text_frame, source_language, target_language,
                                use_reasoning, temperature, max_tokens, top_p
                            )
                            
                            logger.info(f"Translated speaker notes with preserved formatting")
                        
                        processed_items += 1
                        pbar.update(1)
            
            # Save the translated presentation
            prs.save(output_file)
            logger.info(f"Translated presentation saved to {output_file}")
            return True
            
        except Exception as e:
            logger.error(f"File translation failed: {e}")
            return False
    
    def translate_batch(self, texts, source_language="auto (en-US)", target_language="zh-TW", 
                        use_reasoning=False, temperature=0.7, max_tokens=3000, top_p=0.9):
        """
        Translate a batch of texts with progress tracking.
        
        Args:
            texts (list): List of texts to translate.
            source_language (str): Source language, default is 'auto (en-US)'.
            target_language (str): Target language, default is 'zh-TW'.
            use_reasoning (bool): Whether to enable extended reasoning for improved accuracy.
            temperature (float): Temperature for model creativity, default is 0.7.
            max_tokens (int): Maximum tokens for response, default is 3000.
            top_p (float): Top P for nucleus sampling, default is 0.9.
            
        Returns:
            list: List of translated texts.
        """
        results = []
        for text in tqdm(texts, desc="Batch translation progress"):
            if text.strip():
                translated = self.translate(
                    text, source_language, target_language, 
                    use_reasoning, temperature, max_tokens, top_p
                )
                results.append(translated)
            else:
                results.append("")
        logger.info(f"Completed batch translation of {len(texts)} texts")
        return results
    
    def download_from_s3(self, bucket, key, local_path):
        """
        Download a file from S3 to a local path.
        
        Args:
            bucket (str): S3 bucket name.
            key (str): S3 object key.
            local_path (str): Local path to save the file.
            
        Returns:
            bool: True if download is successful, False otherwise.
        """
        try:
            self.s3_client.download_file(bucket, key, local_path)
            logger.info(f"Downloaded file from S3: {bucket}/{key} to {local_path}")
            return True
        except Exception as e:
            logger.error(f"Error downloading from S3: {e}")
            return False
    
    def upload_to_s3(self, local_path, bucket, key):
        """
        Upload a file to S3 from a local path.
        
        Args:
            local_path (str): Local path of the file to upload.
            bucket (str): S3 bucket name.
            key (str): S3 object key.
            
        Returns:
            bool: True if upload is successful, False otherwise.
        """
        try:
            self.s3_client.upload_file(local_path, bucket, key)
            logger.info(f"Uploaded file to S3: {local_path} to {bucket}/{key}")
            return True
        except Exception as e:
            logger.error(f"Error uploading to S3: {e}")
            return False

def lambda_handler(event, context):
    """
    Lambda function handler to process S3 events for translation jobs.
    
    Args:
        event (dict): Lambda event data.
        context (object): Lambda context object.
        
    Returns:
        dict: Response with status code and message.
    """
    logger.info("Lambda event received: %s", json.dumps(event))
    
    # Check if the event is from S3
    if 'Records' in event and len(event['Records']) > 0:
        record = event['Records'][0]
        if 's3' in record:
            bucket = record['s3']['bucket']['name']
            key = record['s3']['object']['key']
            local_input_file_path = '/tmp/input.pptx'
            local_output_file_path = '/tmp/output.pptx'
            translated_file_key = f"translated/{key.split('/')[-1]}"
            translated_bucket_name = bucket  # Use the same bucket for output, or configure a different one if needed
            
            translator = BedrockTranslator()
            if translator.download_from_s3(bucket, key, local_input_file_path):
                if translator.translate_file(local_input_file_path, local_output_file_path, target_language="zh-TW"):
                    if translator.upload_to_s3(local_output_file_path, translated_bucket_name, translated_file_key):
                        return {
                            'statusCode': 200,
                            'body': json.dumps(f"Successfully translated {key} and uploaded to {translated_bucket_name}/{translated_file_key}")
                        }
                    else:
                        return {
                            'statusCode': 500,
                            'body': json.dumps("Failed to upload translated file to S3")
                        }
                else:
                    return {
                        'statusCode': 500,
                        'body': json.dumps("Failed to translate file")
                    }
            else:
                return {
                    'statusCode': 500,
                    'body': json.dumps("Failed to download input file from S3")
                }
    return {
        'statusCode': 400,
        'body': json.dumps("Invalid event format or no S3 event found")
    }

# Removed local testing components to streamline for Lambda deployment
